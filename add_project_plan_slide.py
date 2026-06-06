"""
Thêm 1 slide "Project Plan — Team Roles & Timeline" vào file pptx,
khớp đúng style các slide hiện có (header navy, thẻ màu teal/purple/orange).

Nội dung dựa trên thành viên nhóm + các phase trong PROJECT_SUMMARY.md.

Chạy:
  python add_project_plan_slide.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

SRC = "DAP391m_ViolenceDetection_Report.pptx"

# ── Palette khớp slide hiện có ──────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x1F, 0x3D)   # header bar
NAVY2   = RGBColor(0x16, 0x20, 0x33)   # card nền tối
TEAL    = RGBColor(0x0D, 0x94, 0x88)
TEAL_LT = RGBColor(0x14, 0xB8, 0xA6)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
ORANGE  = RGBColor(0xD9, 0x77, 0x06)
RED     = RGBColor(0xDC, 0x26, 0x26)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAYTX  = RGBColor(0x47, 0x55, 0x69)
LIGHTTX = RGBColor(0xD1, 0xD5, 0xDB)
SUBTX   = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank)


def box(x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE,
        radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w if line_w else 1)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0):
    """runs: list các paragraph; mỗi paragraph là list (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, b, col) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = col
            r.font.name = "Calibri"
    return tb


# ── HEADER BAR ──────────────────────────────────────────────────────────────
box(0, 0, 10, 0.75, fill=NAVY)
text(0.50, 0.10, 7.5, 0.55,
     [[("Project Plan — Team Roles & Timeline", 20, True, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(7.6, 0.10, 2.3, 0.55,
     [[("DAP391m · Capstone", 11, True, TEAL_LT)]],
     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Sub caption dưới header
text(0.50, 0.80, 9.0, 0.30,
     [[("FPT University, Ho Chi Minh City  ·  Supervisor: Le Nhat Tung  ·  "
        "Target: Scopus Q4 conference", 11, False, GRAYTX)]])

# ════════════════════════════════════════════════════════════════════════════
# SECTION A — TEAM & RESPONSIBILITIES (cột trái)
# ════════════════════════════════════════════════════════════════════════════
text(0.50, 1.18, 4.5, 0.30, [[("Team & Responsibilities", 13, True, NAVY)]])
box(0.50, 1.52, 4.55, 0.03, fill=TEAL)

members = [
    (TEAL,   "Hà Việt Hưng", "SE201122 · Team Lead / Corresponding Author",
     "CGM architecture · X3D-S fine-tuning · paper writing"),
    (PURPLE, "Nguyễn Việt Nhân", "SE201082 · Context & Data",
     "Crowd (YOLOv8n) + Lighting streams · EDA & visualization"),
    (ORANGE, "Nguyễn Thái Kiệt", "SE200734 · Modeling & Experiments",
     "Motion optical-flow stream · ablation study E0–E5"),
    (RED,    "Trần Bảo Nguyên", "SE201012 · Evaluation & Transfer",
     "Cross-dataset RLVS pipeline · metrics & result tables"),
]
cy = 1.66
ch = 0.86
for color, name, role, resp in members:
    box(0.50, cy, 4.55, ch, fill=NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.06)
    box(0.50, cy, 0.10, ch, fill=color)  # dải màu trái
    text(0.68, cy + 0.06, 4.30, ch - 0.10, [
        [(name, 12, True, WHITE)],
        [(role, 9.5, True, LIGHTTX)],
        [(resp, 9.5, False, SUBTX)],
    ], space_after=1, line_spacing=1.0)
    cy += ch + 0.10

# ════════════════════════════════════════════════════════════════════════════
# SECTION B — TIMELINE (cột phải, Gantt theo phase)
# ════════════════════════════════════════════════════════════════════════════
TX = 5.30   # x gốc cột phải
TW = 4.40   # bề rộng vùng phải
text(TX, 1.18, TW, 0.30, [[("Timeline & Status (by phase)", 13, True, NAVY)]])
box(TX, 1.52, TW, 0.03, fill=ORANGE)

# Trục thời gian: 8 cột tuần (W1..W8)
bar_x0   = TX + 1.55          # nơi bắt đầu các thanh bar
bar_area = TW - 1.55          # bề rộng vùng vẽ bar
n_weeks  = 8
col_w    = bar_area / n_weeks

# nhãn tuần
for i in range(n_weeks):
    text(bar_x0 + i * col_w, 1.60, col_w, 0.20,
         [[("W%d" % (i + 1), 7.5, False, SUBTX)]], align=PP_ALIGN.CENTER)

# (phase_label, week_start(0-based), week_span, color, status)
phases = [
    ("Phase 0–1: Data prep & split",      0, 1, TEAL,   "Done"),
    ("Phase 2: X3D-S fine-tune",          1, 1, TEAL,   "Done"),
    ("Phase 3: Context extraction",       2, 2, TEAL,   "Done"),
    ("Phase 4: CGM + ablation E0–E5",     3, 2, TEAL,   "Done"),
    ("Phase 5–6: Cross-dataset (RLVS)",   5, 2, ORANGE, "In Progress"),
    ("Paper writing + Q4 submission",     6, 2, GRAY,   "Pending"),
]
ry = 1.86
rh = 0.40
for label, ws, span, color, status in phases:
    # tên phase (bên trái thanh)
    text(TX, ry + 0.02, 1.52, rh, [[(label, 8.5, True, NAVY)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    # nền track
    box(bar_x0, ry + 0.07, bar_area, rh - 0.16, fill=RGBColor(0xE5, 0xE7, 0xEB),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    # thanh tiến độ
    bx = bar_x0 + ws * col_w
    bw = max(span * col_w - 0.04, 0.2)
    bar = box(bx, ry + 0.07, bw, rh - 0.16, fill=color,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    # nhãn status trên thanh
    text(bx, ry + 0.05, bw, rh - 0.12, [[(status, 7.5, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ry += rh + 0.045

# Chú thích trạng thái (legend) dưới timeline
ly = ry + 0.02
legend = [("Done", TEAL), ("In Progress", ORANGE), ("Pending", GRAY)]
lx = TX
for name, col in legend:
    box(lx, ly + 0.04, 0.16, 0.16, fill=col,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
    text(lx + 0.20, ly - 0.02, 1.2, 0.26, [[(name, 8.5, False, GRAYTX)]])
    lx += 1.45

# ── Footer note ─────────────────────────────────────────────────────────────
box(0.50, 5.32, 9.0, 0.025, fill=RGBColor(0xE5, 0xE7, 0xEB))
text(0.50, 5.34, 9.0, 0.26,
     [[("Current focus: ", 9, True, ORANGE),
       ("Phase 5–6 cross-dataset transfer to RLVS (zero-shot) → then finalize "
        "Methodology, Experiments & Results for submission.", 9, False, GRAYTX)]])

# ── Di chuyển slide mới lên vị trí thứ 3 (ngay sau Outline) ──────────────────
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new = slides[-1]
xml_slides.remove(new)
xml_slides.insert(2, new)   # 0-based: sau Title(0) và Outline(1)

prs.save(SRC)
print("Done. Đã thêm slide 'Project Plan' vào vị trí #3 của", SRC)
print("Tổng số slide hiện tại:", len(prs.slides._sldIdLst))
